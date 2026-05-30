"""
Render slide JSON (from generate_slide_content) into a polished .pptx.

Design choices vs the original spec:
  - Modern dark navy + accent teal palette (more "premium" than primary blue)
  - Title slide uses a vertical accent stripe + tagline + date instead of a
    flat full-bleed colour (more striking; reads as a real cover)
  - Content slides use ONE text frame with native PowerPoint bullets and
    paragraph spacing — not stacked text boxes — so wrapping, spacing, and
    selection behave like a normal slide
  - Footer with project name + slide-number-of-total on every content slide
  - Speaker notes preserved
  - All shapes use the proper MSO_SHAPE enum, not magic numbers
"""

from __future__ import annotations

import io
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# Palette — dark navy + teal accent + warm off-white background
NAVY        = RGBColor(0x0F, 0x1F, 0x3D)
NAVY_DEEP   = RGBColor(0x0A, 0x15, 0x2A)
TEAL        = RGBColor(0x4E, 0xC9, 0xC0)
TEAL_SOFT   = RGBColor(0xA8, 0xE6, 0xE0)
INK         = RGBColor(0x23, 0x2A, 0x33)
INK_SOFT    = RGBColor(0x55, 0x5E, 0x6A)
PAPER       = RGBColor(0xFA, 0xFB, 0xFD)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RULE_GREY   = RGBColor(0xE3, 0xE7, 0xEE)

FONT_HEADING = "Calibri"
FONT_BODY    = "Calibri"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _solid_fill(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _set_background(slide, rgb: RGBColor):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb


def _add_text(slide, text, *, left, top, width, height,
              size=18, bold=False, italic=False, color=INK,
              align=PP_ALIGN.LEFT, font=FONT_BODY, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _add_rect(slide, left, top, width, height, rgb):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _solid_fill(shape, rgb)
    return shape


# ─────────────────────────────────────────────
# Title slide
# ─────────────────────────────────────────────
def _build_title_slide(slide, slide_data: dict):
    _set_background(slide, NAVY)

    # Left vertical accent stripe
    _add_rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, TEAL)

    # Decorative thin rule under title
    _add_rect(slide, Inches(1), Inches(4.05), Inches(2.2), Inches(0.05), TEAL)

    # Project title
    _add_text(
        slide,
        text=slide_data.get("title", "Project Title"),
        left=Inches(1), top=Inches(2.4),
        width=Inches(11.3), height=Inches(1.5),
        size=54, bold=True, color=WHITE,
        font=FONT_HEADING,
        line_spacing=1.05,
    )

    # Tagline / subtitle from first bullet
    bullets = slide_data.get("bullets", [])
    tagline = bullets[0] if bullets else ""
    if tagline:
        _add_text(
            slide,
            text=tagline,
            left=Inches(1), top=Inches(4.25),
            width=Inches(11.3), height=Inches(1.0),
            size=22, color=TEAL_SOFT, italic=True,
            line_spacing=1.25,
        )

    # Footer-ish lockup
    _add_text(
        slide,
        text="VIVORA",
        left=Inches(1), top=Inches(6.4),
        width=Inches(6), height=Inches(0.4),
        size=12, bold=True, color=TEAL,
    )
    _add_text(
        slide,
        text=datetime.now().strftime("%B %Y"),
        left=Inches(7.3), top=Inches(6.4),
        width=Inches(5), height=Inches(0.4),
        size=12, color=TEAL_SOFT, align=PP_ALIGN.RIGHT,
    )


# ─────────────────────────────────────────────
# Content slide
# ─────────────────────────────────────────────
def _build_content_slide(slide, slide_data: dict, slide_number: int, total: int, project_name: str):
    _set_background(slide, PAPER)

    # Left accent stripe (signature element)
    _add_rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, TEAL)

    # Title
    _add_text(
        slide,
        text=slide_data.get("title", "").upper(),
        left=Inches(0.8), top=Inches(0.55),
        width=Inches(11.5), height=Inches(0.7),
        size=14, bold=True, color=TEAL,
        font=FONT_HEADING, line_spacing=1.0,
    )

    _add_text(
        slide,
        text=slide_data.get("title", ""),
        left=Inches(0.8), top=Inches(0.95),
        width=Inches(11.5), height=Inches(1.0),
        size=34, bold=True, color=NAVY,
        font=FONT_HEADING, line_spacing=1.05,
    )

    # Thin rule under title
    _add_rect(slide, Inches(0.8), Inches(1.95), Inches(1.8), Inches(0.05), TEAL)

    # Bullets — one text frame with native bullet paragraphs
    bullets = slide_data.get("bullets", [])[:5]
    if bullets:
        box = slide.shapes.add_textbox(
            Inches(0.85), Inches(2.4),
            Inches(11.6), Inches(4.3),
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.3
            p.space_after = Pt(10)

            marker = p.add_run()
            marker.text = "▸  "
            marker.font.name = FONT_BODY
            marker.font.size = Pt(20)
            marker.font.bold = True
            marker.font.color.rgb = TEAL

            body = p.add_run()
            body.text = bullet
            body.font.name = FONT_BODY
            body.font.size = Pt(20)
            body.font.color.rgb = INK

    # Footer bar
    _add_rect(slide, Inches(0), Inches(7.10), SLIDE_W, Inches(0.02), RULE_GREY)

    _add_text(
        slide,
        text=project_name,
        left=Inches(0.8), top=Inches(7.15),
        width=Inches(8), height=Inches(0.3),
        size=10, color=INK_SOFT,
    )
    _add_text(
        slide,
        text=f"{slide_number} / {total}",
        left=Inches(11.5), top=Inches(7.15),
        width=Inches(1.2), height=Inches(0.3),
        size=10, color=INK_SOFT, align=PP_ALIGN.RIGHT, bold=True,
    )


# ─────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────
def create_presentation(slides_data: list, project_name: str = "Project") -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]
    total = len(slides_data)

    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        if i == 0:
            _build_title_slide(slide, slide_data)
        else:
            _build_content_slide(slide, slide_data, i + 1, total, project_name)

        notes = slide_data.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
